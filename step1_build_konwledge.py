#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
批量调用 MinerU 解析指定目录下的所有文件，并根据解析出的所有 MD 文件构建知识库。

功能：
1. 遍历输入目录，使用 MinerU 解析每个文件，输出到指定目录。
   - 支持跳过已解析文件（可被 --force-reparse 覆盖）。
2. 解析完成后，扫描输出目录，收集所有生成的 .md 文件。
3. **智能切分**：根据文件所在的子目录（'h1' 或 'newline'）自动选择切分方法。
   - 位于 'h1' 子目录下的文件，按一级标题（"# "）切分。
   - 位于 'newline' 子目录下的文件，按空行切分。
   - 其他文件使用默认的切分方法。
4. 过滤掉 token 数小于指定阈值的分片。
5. 将所有合格的分片存入知识库的 kv_store_text_chunks.json 文件。
6. 仅生成文本分片 JSON；向量化与 Chroma 写入由 step2.5_json2chroma.py 完成。
"""

# ... (从 # export http_proxy... 到 class MineruParser 之前的所有 import 和类定义保持不变) ...
import os
import sys
import json
import time
import logging
import subprocess
import tempfile
import zlib
import base64
import struct
import re
import shutil
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union
from queue import Queue, Empty
import threading
import argparse

# export http_proxy=127.0.0.1:7890
# export https_proxy=127.0.0.1:7890
# ---------------------------
# 异常类型
# ---------------------------
class MineruExecutionError(Exception):
    def __init__(self, return_code: int, errors: List[str]):
        self.return_code = return_code
        self.errors = errors or []
        super().__init__(self.__str__())

    def __str__(self):
        msg = f"MinerU 执行失败，返回码={self.return_code}"
        if self.errors:
            msg += f"，错误信息（截断）：{self.errors[:3]}"
        return msg


# ---------------------------
# MinerU 解析器（独立封装）
# ---------------------------
class MineruParser:
    """独立的 MinerU 2.0 文档解析器封装"""

    logger = logging.getLogger(__name__)

    IMAGE_FORMATS = {".png", ".jpeg", ".jpg", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
    OFFICE_FORMATS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"}
    TEXT_FORMATS = {".txt", ".md"}

    def __init__(self) -> None:
        pass

    @staticmethod
    def check_installation() -> bool:
        """检查 mineru 是否安装可用"""
        try:
            result = subprocess.run(
                ["mineru", "--version"],
                capture_output=True, text=True, check=True, encoding="utf-8", errors="ignore"
            )
            logging.debug(f"MinerU 版本：{result.stdout.strip()}")
            return True
        except (subprocess.CalledProcessError, FileNotFoundError):
            logging.error("未检测到 MinerU 2.0。请先安装：pip install -U 'mineru[core]'")
            return False

    @staticmethod
    def _run_mineru_command(cmd: List[str]) -> None:
        """运行 mineru 命令行工具并实时捕获输出"""
        logging.info(f"执行 MinerU 命令：{' '.join(cmd)}")
        process = subprocess.Popen(
            cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True,
            encoding="utf-8", errors="ignore", bufsize=1
        )
        stdout_queue, stderr_queue = Queue(), Queue()

        def enqueue_output(pipe, queue, prefix):
            try:
                for line in iter(pipe.readline, ""):
                    if line.strip():
                        queue.put((prefix, line.strip()))
                pipe.close()
            except Exception as e:
                queue.put((prefix, f"读取{prefix}错误：{e}"))

        threading.Thread(target=enqueue_output, args=(process.stdout, stdout_queue, "STDOUT"), daemon=True).start()
        threading.Thread(target=enqueue_output, args=(process.stderr, stderr_queue, "STDERR"), daemon=True).start()

        error_lines = []
        while process.poll() is None:
            for q, log_func, err_list in [(stdout_queue, logging.info, None), (stderr_queue, logging.warning, error_lines)]:
                try:
                    while True:
                        _, line = q.get_nowait()
                        if "error" in line.lower():
                            logging.error(f"[MinerU] {line}")
                            if err_list is not None: err_list.append(line)
                        else:
                            log_func(f"[MinerU] {line}")
                except Empty:
                    pass
            time.sleep(0.1)
        
        process.wait()
        # 清空剩余日志
        for q, log_func, err_list in [(stdout_queue, logging.info, None), (stderr_queue, logging.warning, error_lines)]:
            try:
                while True:
                    _, line = q.get_nowait()
                    if "error" in line.lower():
                        logging.error(f"[MinerU] {line}")
                        if err_list is not None: err_list.append(line)
                    else:
                        log_func(f"[MinerU] {line}")
            except Empty:
                pass

        if process.returncode != 0 or error_lines:
            raise MineruExecutionError(process.returncode, error_lines)

    def parse_document(
        self, file_path: Path, output_dir: Path, **kwargs
    ) -> None:
        """统一的文档解析入口"""
        ext = file_path.suffix.lower()
        method = kwargs.get("method", "auto")

        # 对 Office/Text 文件，先转为 PDF（Office 可走纯文本回退，见文末 except）
        temp_pdf = None
        input_to_mineru = file_path
        
        if ext in self.OFFICE_FORMATS:
            logging.info(f"检测到 Office 文件 {ext}，将先通过 LibreOffice 转为 PDF。")
            resolved = self._office_resolve_pdf_or_text_fallback(file_path, output_dir, kwargs.get("method", "auto"))
            if resolved is None:
                # 已写入降级 .md（无 MinerU / 或无 PDF），本条 Office 预处理结束。
                logging.info(f"✓ Office 文件已仅以文本回落入库路径: {file_path.name}")
                return
            temp_pdf = resolved
            input_to_mineru = temp_pdf
            method = "auto"
        elif ext in self.TEXT_FORMATS:
            logging.info(f"检测到文本文件 {ext}，将先转换为 PDF。")
            temp_pdf = self._convert_to_pdf(file_path, output_dir, "text")
            input_to_mineru = temp_pdf
            method = "auto"
        elif ext in self.IMAGE_FORMATS:
            method = "ocr" # 图片强制使用 ocr
        
        cmd = ["mineru", "-p", str(input_to_mineru), "-o", str(output_dir), "-m", method]
        # 添加其他可选参数
        for key, val in kwargs.items():
            if key in ["method"]: continue
            if val is not None:
                if isinstance(val, bool):
                    if val: cmd.append(f"--{key.replace('_', '-')}")
                else:
                    cmd.extend([f"--{key.replace('_', '-')}", str(val)])
        
        try:
            self._run_mineru_command(cmd)
        except Exception as mineru_err:
            # LibreOffice → PDF OK，但 MinerU 解析失败 → 再走纯文本回退
            if ext in self.OFFICE_FORMATS:
                fb = self._write_office_text_fallback_md(
                    file_path, output_dir, method=kwargs.get("method", "auto")
                )
                if fb is not None:
                    logging.warning(
                        f"[Office 回退] MinerU 解析失败 ({mineru_err!r}), "
                        f"已写入降级 Markdown（无版面/插图）: {fb}"
                    )
                    return
            raise
        finally:
            if temp_pdf and temp_pdf.exists():
                try:
                    temp_pdf.unlink() # 清理临时 PDF
                except OSError as e:
                    logging.warning(f"无法删除临时 PDF 文件: {temp_pdf}, {e}")

    def _office_resolve_pdf_or_text_fallback(
        self,
        file_path: Path,
        output_dir: Path,
        mineru_method: str,
    ) -> Optional[Path]:
        """Office → 临时 PDF 供 MinerU；若 LibreOffice 完全失败但有 python-docx/pptx 则写入降级 md 并返回 None。

        Returns:
            Path — 可用的临时 PDF；
            None — 已仅用纯文本Markdown回退，`parse_document` 应跳过后续 MinerU。
        """
        file_path = file_path.expanduser().resolve()
        if not file_path.is_file():
            raise FileNotFoundError(f"源 Office 文件不存在: {file_path}")

        pdf_dir = output_dir / f"{file_path.stem}_temp_pdf"
        shutil.rmtree(pdf_dir, ignore_errors=True)
        pdf_dir.mkdir(parents=True, exist_ok=True)
        expected_pdf = pdf_dir / f"{file_path.stem}.pdf"

        profile = tempfile.mkdtemp(prefix="lo_user_")
        profile_uri = Path(profile).resolve().as_uri()

        candidates: List[str] = []
        env_bin = (os.environ.get("LIBREOFFICE_BIN") or "").strip()
        if env_bin:
            candidates.append(env_bin)
        for nm in ("soffice", "libreoffice"):
            p = shutil.which(nm)
            if p and os.path.isfile(p):
                rp = os.path.realpath(p)
                if rp not in {os.path.realpath(c) for c in candidates}:
                    candidates.append(p)

        proc_err: Optional[subprocess.CompletedProcess[str]] = None
        chosen_exe: Optional[str] = None

        try:
            for exe in candidates:
                # 避免多次尝试时残留失败产生的半份 PDF
                for stale in pdf_dir.glob("*.pdf"):
                    try:
                        stale.unlink()
                    except OSError:
                        pass
                cmd = [
                    exe,
                    "--headless",
                    f"-env:UserInstallation={profile_uri}",
                    "--norestore",
                    "--nofirststartwizard",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(pdf_dir.resolve()),
                    str(file_path),
                ]
                proc = subprocess.run(cmd, check=False, capture_output=True, text=True)
                made = sorted(
                    pdf_dir.glob("*.pdf"),
                    key=lambda p: p.stat().st_mtime if p.exists() else 0,
                    reverse=True,
                )
                ok = proc.returncode == 0 and bool(made)
                proc_err = proc
                chosen_exe = exe
                if ok:
                    break
        finally:
            shutil.rmtree(profile, ignore_errors=True)

        pdf_candidates = sorted(
            pdf_dir.glob("*.pdf"),
            key=lambda p: p.stat().st_mtime,
            reverse=True,
        )

        chosen: Optional[Path] = None
        if expected_pdf.exists():
            chosen = expected_pdf
        elif pdf_candidates:
            chosen = pdf_candidates[0]
            logging.info(
                f"LibreOffice 输出文件名与预期不完全一致；选用最新 PDF：{chosen.name}"
            )

        if chosen is not None:
            logging.debug(f"[Lo→PDF] 使用 {chosen_exe}，产出 {chosen}")
            return chosen

        # ── LibreOffice 完全失败 ──尝试纯文本回退（不要求 MinerU / PDF） ──
        stderr_snip = (proc_err.stderr or "")[:800] if proc_err else "(无stderr)"
        logging.error(
            f"LibreOffice 未产出 PDF ({file_path.name})。stderr 片段: {stderr_snip}"
        )
        fb = self._write_office_text_fallback_md(file_path, output_dir, method=mineru_method)
        if fb is not None:
            logging.warning(
                f"已降级为纯文本 Markdown（请 pip install LibreOffice / 或使用 docker-full 以获得 PDF+MinerU 完整能力）→ {fb}"
            )
            return None

        stdout_snip = (proc_err.stdout or "")[:400] if proc_err else ""
        tried = ",".join(candidates) if candidates else "(未找到 libreoffice/soffice)"
        raise RuntimeError(
            f"LibreOffice→PDF 失败且纯文本回退不可用（请先安装 LibreOffice **或** `pip install python-docx python-pptx`）。"
            f"\n尝试过: [{tried}]\nstdout: {stdout_snip}"
        )

    def _write_office_text_fallback_md(
        self, file_path: Path, output_dir: Path, method: str
    ) -> Optional[Path]:
        """MinerU/LibreOffice 均失败时：用 python-docx / python-pptx 直接抽纯文本，
        写到与 MinerU 一致的目录结构 `{stem}/{method}/{stem}.md`，供后续分片入库。

        Returns:
            md 路径，若依赖缺失或格式不支持则为 None。"""
        ext = file_path.suffix.lower()
        stem = file_path.stem
        chunks: List[str] = []

        if ext == ".docx":
            try:
                from docx import Document  # python-docx
            except ImportError:
                logging.error("纯文本回退需要: pip install python-docx")
                return None
            doc = Document(str(file_path.resolve()))
            for para in doc.paragraphs:
                t = (para.text or "").strip()
                if t:
                    chunks.append(t)
            for tbl in doc.tables:
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        chunks.append(" | ".join(cells))

        elif ext == ".pptx":
            try:
                from pptx import Presentation  # python-pptx
            except ImportError:
                logging.error("纯文本回退需要: pip install python-pptx")
                return None
            prs = Presentation(str(file_path.resolve()))
            for i, slide in enumerate(prs.slides, start=1):
                lines: List[str] = []
                if slide.shapes:
                    for shape in slide.shapes:
                        try:
                            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                                txt = shape.text.strip()
                                if txt:
                                    lines.append(txt.replace("\x0b", " "))
                        except Exception:
                            continue
                if lines:
                    chunks.append(f"# 幻灯片 {i}\n" + "\n".join(lines))

        elif ext == ".ppt":
            logging.warning("[Office 回退] 旧版 .ppt 无法用 python-pptx 读取，请先转 .pptx 或安装 LibreOffice。")
            return None
        else:
            logging.warning(f"[Office 回退] 暂不支持的格式: {ext}")
            return None

        combined = (
            f"# {stem}\n\n"
            "_（本内容由 python-docx / python-pptx 纯文本回退生成：无 OCR、无图示路径、无原版式。建议修复 LibreOffice+MinerU 后重新入库。）_\n\n"
            + "\n\n".join(chunks).strip()
        )
        if not combined.strip():
            return None

        target_dir = output_dir / stem / method
        target_dir.mkdir(parents=True, exist_ok=True)
        md_out = target_dir / f"{stem}.md"
        md_out.write_text(combined + "\n", encoding="utf-8")
        return md_out.resolve()

    def _convert_to_pdf(self, file_path: Path, output_dir: Path, file_type: str) -> Path:
        """将 **纯文本** 文件转为 PDF。（Office/PPT/DOC 走 `_office_resolve_pdf_or_text_fallback`。）"""
        if file_type != "text":
            raise ValueError(f"_convert_to_pdf 仅接受 file_type='text'，收到: {file_type!r}")
        pdf_dir = output_dir / f"{file_path.stem}_temp_pdf"
        pdf_dir.mkdir(exist_ok=True)
        pdf_path = pdf_dir / f"{file_path.stem}.pdf"

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            c = canvas.Canvas(str(pdf_path), pagesize=A4)
            textobject = c.beginText(40, 800)
            with file_path.open("r", encoding="utf-8", errors="ignore") as f:
                for line in f:
                    textobject.textLine(line.rstrip())
            c.drawText(textobject)
            c.save()
        except Exception as e:
            raise RuntimeError(f"使用 reportlab 转换文本失败: {e}")
        
        if not pdf_path.exists():
            raise FileNotFoundError(f"转换为 PDF 后文件未找到: {pdf_path}")
        return pdf_path

# ---------------------------
# 辅助函数
# ---------------------------
def setup_logging():
    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

def get_output_paths(file_path: Path, out_dir_for_dir: Path, method: str) -> Tuple[Optional[Path], Optional[Path]]:
    """获取 MinerU 输出的 md 和 json 文件路径，适配新旧目录结构"""
    stem = file_path.stem
    new_style_dir = out_dir_for_dir / stem / method
    if new_style_dir.is_dir():
        md = new_style_dir / f"{stem}.md"
        json_file = new_style_dir / f"{stem}_content_list.json"
        return (md if md.exists() else None, json_file if json_file.exists() else None)
    
    old_md = out_dir_for_dir / f"{stem}.md"
    old_json = out_dir_for_dir / f"{stem}_content_list.json"
    if old_json.exists():
        return (old_md if old_md.exists() else None, old_json)
        
    return None, None
    
# ---------------------------
# 知识库构建相关函数
# ---------------------------
def split_md_by_h1(md_text: str) -> List[str]:
    """使用正则表达式按一级标题切分 MD 文本，并保留标题"""
    if not md_text or not md_text.strip(): return []
    
    # (?m) 允许多行模式, ^ 匹配行首。() 保留分隔符
    sections = re.split(r'(?m)^#\s', md_text)
    if len(sections) <= 1: return [md_text.strip()]
        
    result = []
    if sections[0].strip(): result.append(sections[0].strip())
    
    for sec in sections[1:]:
        if sec.strip(): result.append(f"# {sec.strip()}")
            
    return result

def split_by_newline(text: str) -> List[str]:
    """使用正则表达式按空行（一个或多个换行）切分文本"""
    if not text or not text.strip(): return []
    
    sections = re.split(r'\n\s*\n', text.strip())
    return [sec.strip() for sec in sections if sec.strip()]

def extract_image_paths_and_clean_markdown(text: str) -> Tuple[List[str], str]:
    """
    提取 Markdown 图片路径并清理图片语法，避免 hash 路径污染向量语义。
    仅处理 markdown 图片语法: ![alt](path)
    """
    if not text:
        return [], ""

    image_paths = re.findall(r'!\[[^\]]*\]\(([^)]+)\)', text)
    cleaned_text = re.sub(r'!\[[^\]]*\]\([^)]+\)', '', text)
    # 清理过多空行
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text).strip()
    return image_paths, cleaned_text

def count_tokens(text: str) -> int:
    """优先使用 tiktoken 计数，否则回退到正则近似"""
    try:
        import tiktoken
        enc = tiktoken.get_encoding("cl100k_base")
        return len(enc.encode(text, disallowed_special=()))
    except (ImportError, Exception):
        return len(re.findall(r'[\u4e00-\u9fa5]|\w+', text))

def make_doc_id(file_name: str) -> str:
    import hashlib
    return "doc-" + hashlib.md5(file_name.encode("utf-8")).hexdigest()

def make_chunk_id(doc_id: str, order_idx: int, content_hash_base: str) -> str:
    import hashlib
    base = f"{doc_id}|{order_idx}|{content_hash_base}"
    return "chunk-" + hashlib.md5(base.encode("utf-8")).hexdigest()

def floats_to_zlib_b64(vec: List[float]) -> str:
    """将 float 列表压缩并编码"""
    if not vec: return ""
    packed = struct.pack(f"<{len(vec)}f", *vec)
    return base64.b64encode(zlib.compress(packed, level=9)).decode("ascii")

def matrix_to_zlib_b64(mat: List[List[float]]) -> str:
    """将向量矩阵压缩并编码"""
    if not mat or not mat[0]: return ""
    flat_list = [item for row in mat for item in row]
    packed = struct.pack(f"<{len(flat_list)}f", *flat_list)
    return base64.b64encode(zlib.compress(packed, level=9)).decode("ascii")

# ---------------------------
# 主流程
# ---------------------------
def main():
    setup_logging()
    parser = argparse.ArgumentParser(description="批量解析文件并根据目录结构智能构建 RAG 知识库")
    parser.add_argument("--input-dir", default="./知识库构建数据/全局RAG", help="输入文件根目录，默认构建 all 库")
    parser.add_argument("--output-dir", default="./rag_output_all_md", help="MinerU 解析输出根目录")
    parser.add_argument("--kb-dir", default="./rag_data/all", help="最终知识库存储根目录")
    
    parser.add_argument("--min-chunk-tokens", type=int, default=25, help="分片最小 token 数，小于此值将被过滤")
    
    # --- 修改：更新了参数说明 ---
    parser.add_argument("--split-method", default="h1", choices=["h1", "newline"],
                        help="默认切分方法，当文件不位于 'h1' 或 'newline' 子目录时使用")
                        
    parser.add_argument("--force-reparse", action="store_true", help="强制重新解析所有文件")
    # MinerU 相关参数
    parser.add_argument("--method", default="auto", choices=["auto", "txt", "ocr"], help="PDF 解析方法；auto 会保留版面/图片等结构信息")
    parser.add_argument("--lang", default=None, help="OCR 语言")
    parser.add_argument("--device", default="cuda:0", help="推理设备")
    parser.add_argument("--backend", default=None, choices=["pipeline", "vlm-transformers", "vlm-vllm-engine", "vlm-http-client"], help="MinerU 后端")
    parser.add_argument("--source", default="modelscope", choices=["huggingface", "modelscope", "local"], help="MinerU 模型来源")
    parser.add_argument("--vram", type=int, default=None, help="单进程占用 GPU 显存上限")
    parser.add_argument("--formula", default=None, choices=["true", "false"], help="是否启用公式解析")
    parser.add_argument("--table", default=None, choices=["true", "false"], help="是否启用表格解析")
    args = parser.parse_args()

    input_root, output_root, kb_root = Path(args.input_dir), Path(args.output_dir), Path(args.kb_dir)
    kv_store_path = kb_root / "kv_store_text_chunks.json"
    kb_images_root = kb_root / "images"
    
    for p in [input_root, output_root, kb_root]: p.mkdir(parents=True, exist_ok=True)

    logging.info(f"输入目录: {input_root}\n解析输出目录: {output_root}\n知识库目录: {kb_root}")
    
    mineru = MineruParser()
    if not mineru.check_installation(): sys.exit(1)

    start_ts = time.time()
    
    # --- 解析阶段 (此阶段逻辑不变，它会保留输入目录的结构) ---
    all_files = [p for p in input_root.rglob("*") if p.is_file()]
    logging.info(f"--- 开始解析阶段，发现 {len(all_files)} 个文件 ---")
    for i, fpath in enumerate(all_files):
        logging.info(f"[{i+1}/{len(all_files)}] 处理: {fpath.relative_to(input_root)}")
        # MinerU 会将输出保存在与输入结构对应的子目录中
        rel_dir = fpath.parent.relative_to(input_root)
        out_dir_for_dir = output_root / rel_dir
        out_dir_for_dir.mkdir(parents=True, exist_ok=True)

        _, json_path = get_output_paths(fpath, out_dir_for_dir, args.method)

        if not args.force_reparse and json_path and json_path.exists():
            logging.info("  -> 检测到已解析结果，跳过解析。")
        else:
            try:
                mineru.parse_document(
                    fpath,
                    out_dir_for_dir,
                    method=args.method,
                    lang=args.lang,
                    device=args.device,
                    backend=args.backend,
                    source=args.source,
                    vram=args.vram,
                    formula=args.formula,
                    table=args.table,
                )
                logging.info(f"  -> 解析成功: {fpath.name}")
            except (MineruExecutionError, RuntimeError) as e:
                logging.error(f"  -> 解析失败: {fpath.name}, 错误: {e}")
    logging.info(f"--- 解析阶段完成，耗时: {time.time() - start_ts:.2f}s ---")

    # --- 知识库构建阶段 ---
    logging.info(f"--- 开始知识库构建阶段 (根据目录自动选择切分方法, 最小 token 数: {args.min_chunk_tokens}) ---")
    
    all_md_files = list(output_root.rglob("*.md"))
    logging.info(f"在 {output_root} 中找到 {len(all_md_files)} 个 .md 文件。")
    
    kv_store, now_ts = {}, int(time.time())
    for md_path in all_md_files:
        original_filename = md_path.stem
        try:
            md_text = md_path.read_text(encoding="utf-8", errors="ignore")
            
            # --- 核心修改：根据文件路径动态决定切分方法 ---
            relative_md_path = md_path.relative_to(output_root)
            current_split_method = args.split_method  # 默认方法

            if len(relative_md_path.parts) > 1:
                top_level_dir = relative_md_path.parts[0]
                if top_level_dir == 'h1' or top_level_dir == 'newline':
                    current_split_method = top_level_dir
            
            logging.info(f"  -> 使用 '{current_split_method}' 方法切分文件: {original_filename}")

            if current_split_method == 'h1':
                sections = split_md_by_h1(md_text)
            else: # 包含了 'newline' 的情况
                sections = split_by_newline(md_text)
            
            doc_id = make_doc_id(original_filename)

            doc_image_paths_set = set()
            doc_output_root = md_path.parent
            doc_images_src = doc_output_root / "images"
            doc_images_dst = kb_images_root / doc_id

            for idx, content in enumerate(sections):
                image_paths, cleaned_content = extract_image_paths_and_clean_markdown(content)
                # 去重并记录 chunk 引用的图片，后续复制到知识库目录
                chunk_image_paths = []
                for rel_img in image_paths:
                    norm_rel = rel_img.strip()
                    if norm_rel:
                        doc_image_paths_set.add(norm_rel)
                        chunk_image_paths.append(f"images/{doc_id}/{Path(norm_rel).name}")

                if not cleaned_content:
                    continue

                tokens = count_tokens(cleaned_content)
                if tokens < args.min_chunk_tokens:
                    logging.warning(f"  -> 分片过短 ({tokens} tokens)，已跳过。文件: {original_filename}, 内容: '{cleaned_content[:50].replace(chr(10), ' ')}...'")
                    continue
                
                chunk_id = make_chunk_id(doc_id, idx, cleaned_content)
                kv_store[chunk_id] = {
                    "_id": chunk_id, "tokens": tokens, "content": cleaned_content,
                    "chunk_order_index": idx, "full_doc_id": doc_id,
                    "file_path": str(relative_md_path), # 存储相对路径以供追溯
                    "image_paths": sorted(set(chunk_image_paths)),
                    "create_time": now_ts, "update_time": now_ts,
                }

            if doc_images_src.exists() and doc_images_src.is_dir():
                doc_images_dst.mkdir(parents=True, exist_ok=True)
                copied_count = 0
                missing_count = 0
                for rel_img in sorted(doc_image_paths_set):
                    src_path = doc_output_root / rel_img
                    if src_path.exists() and src_path.is_file():
                        shutil.copy2(src_path, doc_images_dst / src_path.name)
                        copied_count += 1
                    else:
                        missing_count += 1
                logging.info(
                    f"  -> 图片归档完成: doc_id={doc_id}, 已复制 {copied_count} 张"
                    + (f", 缺失 {missing_count} 张" if missing_count else "")
                )
        except Exception as e:
            logging.error(f"处理 MD 文件失败: {md_path}, 错误: {e}")
            
    with kv_store_path.open("w", encoding="utf-8") as f:
        json.dump(kv_store, f, ensure_ascii=False, indent=2)
    logging.info(f"KV 分片存储完成: {kv_store_path} (共 {len(kv_store)} 个合格分片)")

    logging.info("已按配置跳过 step1 向量化输出（不再生成 vdb_chunks.json），请使用 step2.5_json2chroma.py 完成向量写入 Chroma。")
    logging.info(f"--- 所有任务完成，总耗时: {time.time() - start_ts:.2f}s ---")

if __name__ == "__main__":
    main()
